"""
Builds the "How the Repair Loop Works" deck as a real .pptx file —
same 17-slide structure and verified numbers as the HTML version, native
PowerPoint shapes/connectors/charts throughout (nothing is a screenshot).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.oxml.ns import qn
import copy

# ── palette ──────────────────────────────────────────────────────────────
INK      = RGBColor(0x1B, 0x21, 0x1F)
MUTED    = RGBColor(0x5C, 0x6B, 0x66)
FAINT    = RGBColor(0x8A, 0x96, 0x8F)
LINE     = RGBColor(0xDC, 0xE1, 0xDD)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xF6, 0xF7, 0xF5)
ACCENT   = RGBColor(0x1F, 0x6F, 0x5C)
ACCENT_SOFT = RGBColor(0xE4, 0xF1, 0xEC)
RESOLVED = RGBColor(0x2E, 0x7D, 0x4F)
RESOLVED_BG = RGBColor(0xE4, 0xF1, 0xE7)
GAVEUP   = RGBColor(0xB2, 0x3A, 0x2E)
GAVEUP_BG = RGBColor(0xFB, 0xE8, 0xE5)

MONO = "Consolas"
SANS = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── low-level helpers ───────────────────────────────────────────────────
def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def textbox(slide, left, top, width, height, text, font=SANS, size=14,
            color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def multirun_textbox(slide, left, top, width, height, runs, size=14,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """runs: list of (text, font, size, color, bold, italic) tuples on one paragraph"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    for text, font, sz, color, bold, italic in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def rounded_box(slide, left, top, width, height, fill=SURFACE, line_color=LINE,
                 line_w=1.0, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def circle(slide, left, top, diameter, fill=ACCENT, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    shp.shadow.inherit = False
    return shp


def add_text_in_shape(shape, text, font=SANS, size=12, color=INK, bold=False,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tf


def connector_arrow(slide, x1, y1, x2, y2, color=FAINT, weight=1.75):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    line = conn.line._get_or_add_ln()
    tail = line.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    line.append(tail)
    return conn


def kicker_and_title(slide, kicker, title, dek=None, title_size=28):
    # kicker dot + text
    circle(slide, Inches(0.7), Inches(0.56), Inches(0.09), fill=ACCENT)
    textbox(slide, Inches(0.86), Inches(0.44), Inches(8), Inches(0.32),
            kicker.upper(), font=MONO, size=11, color=ACCENT, bold=True)
    textbox(slide, Inches(0.7), Inches(0.82), Inches(11.9), Inches(0.9),
            title, font=MONO, size=title_size, color=INK, bold=True, line_spacing=1.05)
    y = Inches(0.82) + Inches(0.1 + title_size * 0.014)
    if dek:
        dek_top = Inches(1.55) if title_size <= 28 else Inches(2.0)
        textbox(slide, Inches(0.7), dek_top, Inches(10.6), Inches(0.9),
                dek, font=SANS, size=14, color=MUTED, line_spacing=1.25)
    return


def page_footer(slide, n, total):
    textbox(slide, Inches(12.55), Inches(7.14), Inches(0.7), Inches(0.3),
            f"{n} / {total}", font=MONO, size=9, color=FAINT, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, left, top, width, item_h=0.62, size=14, gap=0.14):
    """items: list of (lead_bold, rest_text) or (lead_bold, rest_text, sub_text)"""
    y = top
    for item in items:
        lead, rest = item[0], item[1]
        sub = item[2] if len(item) > 2 else None
        badge = rounded_box(slide, left, y, Inches(0.32), Inches(0.32),
                             fill=ACCENT_SOFT, line_color=LINE, line_w=0.75, radius=0.3)
        add_text_in_shape(badge, "→", font=MONO, size=13, color=ACCENT, bold=True)
        runs = [(lead + "  ", SANS, size, ACCENT, True, False),
                (rest, SANS, size, INK, False, False)]
        multirun_textbox(slide, left + Inches(0.48), y - Inches(0.02), width - Inches(0.48),
                          Inches(0.4), runs, size=size, line_spacing=1.15)
        if sub:
            textbox(slide, left + Inches(0.48), y + Inches(0.34), width - Inches(0.48),
                    Inches(0.4), sub, font=SANS, size=size - 2.5, color=MUTED, line_spacing=1.2)
            y += Inches(item_h + 0.18)
        else:
            y += Inches(item_h - 0.15)
        y += Inches(gap)


def set_chart_font(chart, font=SANS, size=11, color=MUTED):
    try:
        chart.font.name = font
        chart.font.size = Pt(size)
        chart.font.color.rgb = color
    except Exception:
        pass


TOTAL_SLIDES = 17

# =========================================================================
# 1 — Title
# =========================================================================
s = new_slide()
circle(s, Inches(0.7), Inches(1.66), Inches(0.1), fill=ACCENT)
textbox(s, Inches(0.88), Inches(1.54), Inches(6), Inches(0.32),
        "A WALKTHROUGH", font=MONO, size=12, color=ACCENT, bold=True)
textbox(s, Inches(0.68), Inches(1.95), Inches(11.5), Inches(1.9),
        "How the Repair Loop\nActually Works", font=MONO, size=44, color=INK,
        bold=True, line_spacing=1.05)
textbox(s, Inches(0.7), Inches(3.55), Inches(9.6), Inches(1.0),
        "What happens when a household robot's plan breaks mid-task — how the system "
        "traces the real cause, why it asks a language model for judgment instead of "
        "solving it by rule alone, and how well the whole thing performs.",
        font=SANS, size=15.5, color=MUTED, line_spacing=1.3)

stats = [("211", "FAILURES DIAGNOSED"), ("91.9%", "FIXED AUTOMATICALLY"), ("392", "TASKS TESTED")]
x = Inches(0.7)
for num, lab in stats:
    textbox(s, x, Inches(4.85), Inches(2.6), Inches(0.6), num, font=MONO, size=28,
            color=ACCENT, bold=True)
    textbox(s, x, Inches(5.45), Inches(2.6), Inches(0.4), lab, font=MONO, size=10.5,
            color=FAINT, bold=True)
    x += Inches(2.35)

# =========================================================================
# 2 — The problem
# =========================================================================
s = new_slide()
kicker_and_title(s, "01 · The problem", "The first plan is rarely the last one",
                  "A language model's first attempt at a plan usually looks reasonable "
                  "— and usually breaks somewhere in the middle. Here, the robot is "
                  "told to wash a load of laundry:", title_size=27)

chips = ["Walk to washing\nmachine  ✓", "Open it  ✓", "Pick up soap  ✓",
         "Pick up a dress  ✓", "Pick up pants  ✗"]
cw, ch, gap = Inches(2.05), Inches(0.85), Inches(0.28)
x = Inches(0.7); y = Inches(2.75)
for i, chip in enumerate(chips):
    fail = (i == len(chips) - 1)
    box = rounded_box(s, x, y, cw, ch,
                       fill=GAVEUP_BG if fail else SURFACE,
                       line_color=GAVEUP if fail else LINE, line_w=1.5 if fail else 1.0,
                       radius=0.16)
    add_text_in_shape(box, chip, font=SANS, size=12.5, color=GAVEUP if fail else INK,
                       bold=fail)
    if i < len(chips) - 1:
        connector_arrow(s, x + cw, y + ch / 2, x + cw + gap, y + ch / 2, color=FAINT, weight=1.5)
    x += cw + gap

textbox(s, Inches(0.7), Inches(4.05), Inches(10.8), Inches(0.6),
        "Both hands are already full. The plan didn't account for that — and without "
        "some way to recover, this is where the task simply ends.",
        font=SANS, size=14.5, color=MUTED, line_spacing=1.3)

runs = [("Even given three separate tries", SANS, 15, INK, True, False),
        (" with no diagnosis or repair — just plan, attempt, and try fresh on failure — "
         "that only succeeds ", SANS, 15, MUTED, False, False),
        ("15–78%", MONO, 15, INK, True, False),
        (" of the time, depending how demanding the task is. Everything from here is what "
         "closes that gap in a single, informed attempt.", SANS, 15, MUTED, False, False)]
multirun_textbox(s, Inches(0.7), Inches(4.75), Inches(11.0), Inches(1.3), runs,
                  size=15, line_spacing=1.35)
page_footer(s, 2, TOTAL_SLIDES)

# =========================================================================
# 3 — Pipeline overview
# =========================================================================
s = new_slide()
kicker_and_title(s, "02 · Overview", "Four steps, every time something breaks",
                  title_size=27)

steps = [("1", "Diagnose", "figure out what kind of\nfailure this is"),
         ("2", "Trace the cause", "find the exact moment\nthings went wrong"),
         ("3", "Search for a fix", "explore ways to satisfy\nwhat's missing"),
         ("4", "Weave it back in", "rejoin the rest of the\noriginal plan and retry")]
bw, bh = Inches(2.62), Inches(1.85)
gap = Inches(0.35)
x = Inches(0.7); y = Inches(2.35)
for num, t, f in steps:
    box = rounded_box(s, x, y, bw, bh, fill=SURFACE, line_color=LINE, line_w=1.0, radius=0.09)
    badge = circle(s, x + Inches(0.22), y + Inches(0.22), Inches(0.4), fill=ACCENT_SOFT)
    add_text_in_shape(badge, num, font=MONO, size=13, color=ACCENT, bold=True)
    textbox(s, x + Inches(0.22), y + Inches(0.82), bw - Inches(0.4), Inches(0.4),
            t, font=SANS, size=15, color=INK, bold=True)
    textbox(s, x + Inches(0.22), y + Inches(1.22), bw - Inches(0.4), Inches(0.6),
            f, font=SANS, size=11, color=MUTED, line_spacing=1.2)
    if num != "4":
        connector_arrow(s, x + bw, y + bh / 2, x + bw + gap, y + bh / 2, color=FAINT, weight=1.5)
    x += bw + gap

loop_box = rounded_box(s, Inches(0.7), Inches(4.55), Inches(6.0), Inches(0.5),
                        fill=RGBColor(0xEE, 0xF1, 0xEE), line_color=None, radius=0.5)
add_text_in_shape(loop_box, "↺   up to three tries per task — each re-runs the whole plan from the start",
                   font=SANS, size=11.5, color=FAINT, align=PP_ALIGN.LEFT)

textbox(s, Inches(0.7), Inches(5.35), Inches(10.8), Inches(0.9),
        "A language model is only asked for input once in this loop, for a short hint at "
        "step 3. Steps 1, 2 and 4 are handled entirely by rule — the same rules the "
        "simulated house itself runs on.", font=SANS, size=14, color=MUTED, line_spacing=1.3)
page_footer(s, 3, TOTAL_SLIDES)

# =========================================================================
# 4 — Step 1: Diagnose
# =========================================================================
s = new_slide()
kicker_and_title(s, "03 · Step 1", "Diagnose: what kind of failure is this?",
                  "The house doesn't explain why an action failed — so the system replays "
                  "what happened and works it out. Three shapes of failure come up in practice:",
                  title_size=27)
bullet_list(s, [
    ("A step was skipped", "Something the action needed was never done.",
     "e.g. trying to put an item away that was never picked up."),
    ("Steps happened out of order", "Everything needed did happen, just in the wrong sequence.",
     "e.g. both hands filled up before the plan tried to pick up a third item."),
    ("The object can't do that", "The action doesn't apply to this object at all.",
     "e.g. asking the character to lie down on something you can't lie on."),
], Inches(0.7), Inches(2.95), Inches(11.2), item_h=0.55, size=15.5, gap=0.28)
page_footer(s, 4, TOTAL_SLIDES)

# =========================================================================
# 5 — Step 2: root cause + strategies
# =========================================================================
s = new_slide()
kicker_and_title(s, "04 · Step 2", "Tracing back to the real cause",
                  "The system replays the plan from the very beginning and watches for the "
                  "exact moment things stopped being true — not the step that just failed, "
                  "but whatever happened earlier that actually caused it. Then it picks how to respond:",
                  title_size=25)

cards = [
    ("Add the missing\nstep", "One small thing was skipped.", "walk over, stand up, face the object", False),
    ("Patch this action", "The object can't do what was asked.", "try a close variation, same spot", False),
    ("Swap the action", "The step itself was the wrong idea.", "ask for a different one entirely", False),
    ("Already done", "The goal is already true.", "drop the now-pointless step", False),
    ("Rebuild the\nsequence", "Everything else — more than one thing needs to change, "
     "or no single-step patch applies.", "hands off to the full search, next", True),
]
cw = Inches(2.18); ch = Inches(2.55); gap = Inches(0.14)
x = Inches(0.7); y = Inches(3.15)
for name, desc, trig, emph in cards:
    box = rounded_box(s, x, y, cw, ch,
                       fill=ACCENT_SOFT if emph else SURFACE,
                       line_color=ACCENT if emph else LINE, line_w=1.5 if emph else 1.0,
                       radius=0.08)
    textbox(s, x + Inches(0.16), y + Inches(0.16), cw - Inches(0.32), Inches(0.7),
            name, font=SANS, size=13.5, color=ACCENT, bold=True, line_spacing=1.1)
    textbox(s, x + Inches(0.16), y + Inches(0.95), cw - Inches(0.32), Inches(1.1),
            desc, font=SANS, size=10.5, color=MUTED, line_spacing=1.25)
    textbox(s, x + Inches(0.16), y + ch - Inches(0.55), cw - Inches(0.32), Inches(0.45),
            trig, font=SANS, size=9, color=FAINT, italic=True, line_spacing=1.15)
    x += cw + gap
page_footer(s, 5, TOTAL_SLIDES)

# =========================================================================
# 6 — Step 3: the search
# =========================================================================
s = new_slide()
kicker_and_title(s, "05 · Step 3", "Searching for a fix",
                  "For the general case, the system explores possible next actions and keeps "
                  "the shortest sequence that actually satisfies what's missing — checked "
                  "directly against the same rules the house itself enforces. Three sources of ideas:",
                  title_size=27)
bullet_list(s, [
    ("Rule-based fixes", "Open a closed door, free up a hand, walk to the right spot.",
     "Worked out directly from what's known about the objects involved — no guessing required."),
    ("A hint from the language model", "A short list of suggested actions naming what might fix it.",
     "Treated as ideas to test, not something accepted on faith."),
    ("What the plan was already trying to do", "Keeps the search from losing track of steps it still needs to cover."),
], Inches(0.7), Inches(3.05), Inches(11.2), item_h=0.55, size=15, gap=0.26)
page_footer(s, 6, TOTAL_SLIDES)

# =========================================================================
# 7 — Step 4: splice + safety nets
# =========================================================================
s = new_slide()
kicker_and_title(s, "06 · Step 4", "Weave it back in, then try again",
                  "The fix gets stitched into the original plan and the whole thing re-runs "
                  "from the top. Three specific safeguards keep that stitching from causing "
                  "new problems:", title_size=27)
bullet_list(s, [
    ("Put things where they belong", "If both hands are full, the fix delivers a held item to where it was actually supposed to end up — not just the floor."),
    ("Reopen what got closed", "If something needed is stuck behind a door another fix just shut, the system reopens it, retrieves what it needs, then restores it."),
    ("Don't discard a necessary step", "A pickup that's the actual cause of a problem still has to happen eventually — it's no longer thrown away just for being the trigger."),
], Inches(0.7), Inches(2.95), Inches(11.2), item_h=0.72, size=15, gap=0.22)
page_footer(s, 7, TOTAL_SLIDES)

# =========================================================================
# 8 — Why not pure search (funnel)
# =========================================================================
s = new_slide()
kicker_and_title(s, "07 · Why a language model", "The search space is enormous, blind",
                  title_size=27)
runs = [("The house understands ", SANS, 14.5, MUTED, False, False),
        ("41", MONO, 14.5, INK, True, False),
        (" distinct actions and holds roughly ", SANS, 14.5, MUTED, False, False),
        ("290", MONO, 14.5, INK, True, False),
        (" objects in an average scene. Checked blindly — before even considering actions "
         "involving two objects at once — that's a lot to get through before finding one valid repair.",
         SANS, 14.5, MUTED, False, False)]
multirun_textbox(s, Inches(0.7), Inches(1.62), Inches(11.4), Inches(0.9), runs,
                  size=14.5, line_spacing=1.3)

# funnel: 3 stages, shrinking
stage_y = Inches(3.15)
# stage 1 (big rect)
b1w, b1h = Inches(2.6), Inches(1.5)
b1 = rounded_box(s, Inches(0.9), stage_y, b1w, b1h, fill=SURFACE, line_color=FAINT, line_w=1.75, radius=0.1)
add_text_in_shape(b1, "~12,000", font=MONO, size=22, color=INK, bold=True)
textbox(s, Inches(0.9), stage_y + b1h + Inches(0.12), b1w, Inches(0.8),
        "raw action–object combinations, checked blindly", font=SANS, size=10.5,
        color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)

connector_arrow(s, Inches(3.5), stage_y + b1h/2, Inches(4.55), stage_y + b1h/2 - Inches(0.25), color=FAINT, weight=1.5)
textbox(s, Inches(3.42), stage_y - Inches(0.42), Inches(1.5), Inches(0.6),
        "language\nmodel\nproposes", font=SANS, size=8.5, color=FAINT, align=PP_ALIGN.CENTER, line_spacing=1.05)

# stage 2 (mid rect)
b2w, b2h = Inches(1.7), Inches(1.05)
b2y = stage_y + Inches(0.22)
b2 = rounded_box(s, Inches(4.65), b2y, b2w, b2h, fill=SURFACE, line_color=ACCENT, line_w=1.75, radius=0.12)
add_text_in_shape(b2, "3", font=MONO, size=22, color=ACCENT, bold=True)
textbox(s, Inches(4.55), b2y + b2h + Inches(0.12), Inches(1.9), Inches(1.0),
        "candidate actions, typically — the median across every failure this session",
        font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)

connector_arrow(s, Inches(6.35), b2y + b2h/2, Inches(7.35), b2y + b2h/2 + Inches(0.2), color=FAINT, weight=1.5)
textbox(s, Inches(6.15), b2y - Inches(0.42), Inches(1.4), Inches(0.4),
        "rules verify", font=SANS, size=8.5, color=FAINT, align=PP_ALIGN.CENTER)

# stage 3 (small circle)
b3d = Inches(0.95)
b3y = b2y + b2h/2 - b3d/2 + Inches(0.2)
b3 = circle(s, Inches(7.45), b3y, b3d, fill=ACCENT)
add_text_in_shape(b3, "1", font=MONO, size=20, color=SURFACE, bold=True)
textbox(s, Inches(7.15), b3y + b3d + Inches(0.12), Inches(1.55), Inches(0.6),
        "repair actually used", font=SANS, size=10, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)

textbox(s, Inches(0.7), Inches(5.95), Inches(11.2), Inches(0.9),
        "Something has to narrow the field before exact checking is even feasible. That "
        "narrowing step is where judgment — not more search — earns its place.",
        font=SANS, size=14, color=MUTED, line_spacing=1.3)
page_footer(s, 8, TOTAL_SLIDES)

# =========================================================================
# 9 — Doctor and lab
# =========================================================================
s = new_slide()
kicker_and_title(s, "08 · Why a language model", "Think of it as a doctor and a lab",
                  title_size=27)

# doctor card
dx, dy, dw, dh = Inches(0.7), Inches(1.95), Inches(5.6), Inches(4.4)
dcard = rounded_box(s, dx, dy, dw, dh, fill=SURFACE, line_color=LINE, line_w=1.0, radius=0.06)
dic = circle(s, dx + Inches(0.25), dy + Inches(0.25), Inches(0.55), fill=ACCENT_SOFT)
add_text_in_shape(dic, "MD", font=SANS, size=13, color=ACCENT, bold=True)
textbox(s, dx + Inches(0.95), dy + Inches(0.35), dw - Inches(1.2), Inches(0.5),
        "The language model is the doctor", font=SANS, size=15, color=INK, bold=True, line_spacing=1.1)
textbox(s, dx + Inches(0.3), dy + Inches(1.15), dw - Inches(0.6), Inches(1.4),
        "It looks at the failure and suggests what might fix it — typically about three "
        "candidate actions per failure, drawn from everyday judgment about how a house "
        "actually works.", font=SANS, size=12.5, color=MUTED, line_spacing=1.35)
textbox(s, dx + Inches(0.3), dy + Inches(2.7), dw - Inches(0.6), Inches(1.4),
        "A doctor's judgment is what narrows down a huge range of possibilities to a short, "
        "plausible list. It isn't the final word.", font=SANS, size=12.5, color=MUTED, line_spacing=1.35)

# lab card
lx = dx + dw + Inches(0.3)
lcard = rounded_box(s, lx, dy, dw, dh, fill=SURFACE, line_color=LINE, line_w=1.0, radius=0.06)
lic = circle(s, lx + Inches(0.25), dy + Inches(0.25), Inches(0.55), fill=RGBColor(0xEE, 0xF1, 0xEE))
add_text_in_shape(lic, "LAB", font=SANS, size=11, color=MUTED, bold=True)
textbox(s, lx + Inches(0.95), dy + Inches(0.35), dw - Inches(1.2), Inches(0.5),
        "The rule-based checker is the lab", font=SANS, size=15, color=INK, bold=True, line_spacing=1.1)
textbox(s, lx + Inches(0.3), dy + Inches(1.15), dw - Inches(0.6), Inches(0.85),
        "It doesn't take the doctor's suggestions on faith. For each one, it can only "
        "answer one of two things:", font=SANS, size=12.5, color=MUTED, line_spacing=1.35)

vpw = (dw - Inches(0.7)) / 2
vp1 = rounded_box(s, lx + Inches(0.3), dy + Inches(2.15), vpw, Inches(1.5), fill=GAVEUP_BG, line_color=None, radius=0.1)
textbox(s, lx + Inches(0.5), dy + Inches(2.35), vpw - Inches(0.4), Inches(0.4),
        "Provably wrong", font=SANS, size=12, color=GAVEUP, bold=True)
textbox(s, lx + Inches(0.5), dy + Inches(2.75), vpw - Inches(0.4), Inches(0.8),
        "rejected outright", font=SANS, size=11, color=GAVEUP, line_spacing=1.25)

vp2 = rounded_box(s, lx + Inches(0.3) + vpw + Inches(0.1), dy + Inches(2.15), vpw, Inches(1.5),
                   fill=RGBColor(0xEE, 0xF1, 0xEE), line_color=None, radius=0.1)
textbox(s, lx + Inches(0.5) + vpw + Inches(0.1), dy + Inches(2.35), vpw - Inches(0.4), Inches(0.4),
        "Can't rule it out", font=SANS, size=12, color=MUTED, bold=True)
textbox(s, lx + Inches(0.5) + vpw + Inches(0.1), dy + Inches(2.75), vpw - Inches(0.4), Inches(0.8),
        "passes through for the search to try", font=SANS, size=11, color=MUTED, line_spacing=1.25)
page_footer(s, 9, TOTAL_SLIDES)

# =========================================================================
# 10 — The lab has final word (guarantee)
# =========================================================================
s = new_slide()
kicker_and_title(s, "09 · Why a language model", "The lab has the final word, always",
                  "A structured search explores only what the lab didn't rule out, and returns "
                  "the shortest sequence that actually resolves the problem. Nothing downstream "
                  "of that search asks the language model anything — it's ordinary, "
                  "deterministic code from there.", title_size=26)

gbox = rounded_box(s, Inches(0.7), Inches(3.3), Inches(11.2), Inches(2.6),
                    fill=ACCENT_SOFT, line_color=ACCENT, line_w=1.75, radius=0.06)
textbox(s, Inches(1.1), Inches(3.65), Inches(10.4), Inches(0.5),
        "The guarantee this gives you", font=SANS, size=17, color=INK, bold=True)
textbox(s, Inches(1.1), Inches(4.3), Inches(10.4), Inches(1.4),
        "The language model cannot make a broken repair pass. At worst, its suggestions "
        "are all rejected — and the step gets dropped, which can cost a goal, but it is "
        "never silently accepted as a fix that doesn't actually hold up. A wrong idea "
        "costs progress. It never costs correctness.",
        font=SANS, size=14.5, color=MUTED, line_spacing=1.4)
page_footer(s, 10, TOTAL_SLIDES)

# =========================================================================
# 11 — The one exception, quantified
# =========================================================================
s = new_slide()
kicker_and_title(s, "10 · Why a language model", "The one honest exception — and how rare it is",
                  "There is exactly one situation where the language model's output goes "
                  "straight into the plan without the lab checking it first: when a step is "
                  "diagnosed as fundamentally the wrong thing to attempt, it's asked to propose "
                  "a direct replacement, accepted as-is.", title_size=24)

runs = [("2", MONO, 62, ACCENT, True, False), (" / 211", MONO, 30, FAINT, True, False)]
multirun_textbox(s, Inches(0.7), Inches(3.55), Inches(5.0), Inches(1.2), runs, anchor=MSO_ANCHOR.BOTTOM)
textbox(s, Inches(4.0), Inches(3.95), Inches(6.5), Inches(0.9),
        "diagnosed failures took this unverified path across both evaluation sets — under 1%.",
        font=SANS, size=13.5, color=MUTED, line_spacing=1.3, anchor=MSO_ANCHOR.BOTTOM)

runs2 = [("Everywhere else — the other ", SANS, 15, MUTED, False, False),
         ("99%+", SANS, 15, INK, True, False),
         (" — every suggestion passes through the lab before anything reaches the plan.",
          SANS, 15, MUTED, False, False)]
multirun_textbox(s, Inches(0.7), Inches(5.1), Inches(10.6), Inches(0.8), runs2, size=15, line_spacing=1.3)
page_footer(s, 11, TOTAL_SLIDES)

# =========================================================================
# 12 — Synthesis
# =========================================================================
s = new_slide()
kicker_and_title(s, "11 · Why a language model", "Judgment narrows it. Verification makes it exact.",
                  title_size=25)
bullet_list(s, [
    ("Judgment picks a short list", "From thousands of raw possibilities down to a handful of plausible ones, in a single step.",
     "No everyday common sense is built into the rules themselves."),
    ("Verification makes it exact", "Nothing reaches the plan unless it genuinely, provably satisfies what's missing."),
    ("Neither half is enough alone", "Judgment without checking would occasionally accept something that looks right but isn't.",
     "Checking without judgment would have to search enormously to find what a person would guess instantly."),
], Inches(0.7), Inches(2.85), Inches(11.4), item_h=0.65, size=15.5, gap=0.3)
page_footer(s, 12, TOTAL_SLIDES)

# =========================================================================
# 13 — Headline results
# =========================================================================
s = new_slide()
kicker_and_title(s, "12 · Results", "Does the loop actually help?", title_size=27)

def headline_card(x, title, noadapt, sda, foot_runs):
    w = Inches(5.55); h = Inches(2.65); y = Inches(2.05)
    box = rounded_box(s, x, y, w, h, fill=SURFACE, line_color=LINE, line_w=1.0, radius=0.06)
    textbox(s, x + Inches(0.3), y + Inches(0.22), w - Inches(0.6), Inches(0.4),
            title.upper(), font=MONO, size=11, color=MUTED, bold=True)
    # bars
    label_w = Inches(2.1); bar_x = x + Inches(0.3) + label_w; bar_w = w - Inches(0.6) - label_w - Inches(0.75)
    by = y + Inches(0.8)
    for lab, val, color in [("No repair, best of 3", noadapt, FAINT), ("With the repair loop", sda, ACCENT)]:
        textbox(s, x + Inches(0.3), by, label_w - Inches(0.1), Inches(0.4), lab, font=SANS, size=11, color=MUTED,
                anchor=MSO_ANCHOR.MIDDLE)
        track = rounded_box(s, bar_x, by + Inches(0.03), bar_w, Inches(0.34), fill=RGBColor(0xEE,0xF1,0xEE),
                             line_color=None, radius=0.3)
        fillw = Emu(int(bar_w * (val / 100.0)))
        if fillw > Emu(int(Inches(0.05))):
            fill = rounded_box(s, bar_x, by + Inches(0.03), fillw, Inches(0.34), fill=color, line_color=None, radius=0.3)
        textbox(s, bar_x + bar_w + Inches(0.08), by, Inches(0.75), Inches(0.4), f"{val}%",
                font=MONO, size=13, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        by += Inches(0.55)
    fbox = rounded_box(s, x + Inches(0.3), y + h - Inches(0.62), w - Inches(0.6), Inches(0.02),
                        fill=LINE, line_color=None, radius=0)
    multirun_textbox(s, x + Inches(0.3), y + h - Inches(0.5), w - Inches(0.6), Inches(0.42),
                      foot_runs, size=11, line_spacing=1.2)

headline_card(Inches(0.7), "Everyday tasks — 342 of them", 77.9, 90.6,
              [("+12.7pt", MONO, 11.5, INK, True, False),
               (" over three independent, unguided attempts.", SANS, 11, FAINT, False, False)])
headline_card(Inches(6.85), "Deliberately hard tasks — 50 of them", 14.6, 50.0,
              [("3.4×", MONO, 11.5, INK, True, False),
               (" better than best-of-3 unguided attempts.", SANS, 11, FAINT, False, False)])

textbox(s, Inches(0.7), Inches(5.05), Inches(11.2), Inches(0.8),
        "The repair loop makes one careful pass with up to three fix attempts — it isn't "
        "“three independent guesses,” it's three informed ones.",
        font=SANS, size=14, color=MUTED, line_spacing=1.3)
page_footer(s, 13, TOTAL_SLIDES)

# =========================================================================
# 14 / 15 — breakdown charts (native pptx charts)
# =========================================================================
STRAT_LABELS = {
    "reconstruct": "Rebuild sequence",
    "insert_prep": "Add missing step",
    "local": "Patch this action",
    "already_satisfied": "Already done",
    "wrong_action": "Swap the action",
}

def strategy_chart_slide(kicker, title, strat_counts, outcome, stats):
    s = new_slide()
    kicker_and_title(s, kicker, title, title_size=27)

    # left: clustered bar of strategy counts
    labels = [STRAT_LABELS[k] for k in strat_counts]
    vals = list(strat_counts.values())
    order = sorted(range(len(vals)), key=lambda i: -vals[i])
    labels = [labels[i] for i in order]
    vals = [vals[i] for i in order]

    textbox(s, Inches(0.7), Inches(1.85), Inches(5.4), Inches(0.35),
            "HOW EACH FAILURE GOT HANDLED", font=MONO, size=11, color=FAINT, bold=True)
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("count", vals)
    gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(2.2),
                                 Inches(5.7), Inches(3.55), cd)
    chart = gframe.chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = True
    dl = chart.plots[0].data_labels
    dl.number_format = "0"; dl.number_format_is_linked = False
    dl.font.size = Pt(11); dl.font.color.rgb = INK; dl.font.name = SANS
    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT
    series.format.line.fill.background()
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.name = SANS
    cat_ax.tick_labels.font.color.rgb = INK
    cat_ax.format.line.color.rgb = LINE
    cat_ax.major_tick_mark = XL_TICK_MARK.NONE
    val_ax = chart.value_axis
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    set_chart_font(chart)

    # right: stacked bar resolved/gave_up
    textbox(s, Inches(6.85), Inches(1.85), Inches(5.4), Inches(0.35),
            "...AND WHETHER IT ACTUALLY WORKED", font=MONO, size=11, color=FAINT, bold=True)
    o_labels = [STRAT_LABELS[k] for k in outcome]
    r_vals = [outcome[k].get("resolved", 0) for k in outcome]
    g_vals = [outcome[k].get("gave_up", 0) for k in outcome]
    totals = [r + g for r, g in zip(r_vals, g_vals)]
    order2 = sorted(range(len(totals)), key=lambda i: -totals[i])
    o_labels = [o_labels[i] for i in order2]
    r_vals = [r_vals[i] for i in order2]
    g_vals = [g_vals[i] for i in order2]

    cd2 = CategoryChartData()
    cd2.categories = o_labels
    cd2.add_series("Fixed", r_vals)
    cd2.add_series("Gave up", g_vals)
    gframe2 = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(6.7), Inches(2.2),
                                  Inches(5.9), Inches(3.55), cd2)
    chart2 = gframe2.chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.legend.include_in_layout = False
    chart2.legend.font.size = Pt(10.5); chart2.legend.font.name = SANS
    plot2 = chart2.plots[0]
    plot2.gap_width = 60
    plot2.series[0].format.fill.solid(); plot2.series[0].format.fill.fore_color.rgb = RESOLVED
    plot2.series[0].format.line.fill.background()
    plot2.series[1].format.fill.solid(); plot2.series[1].format.fill.fore_color.rgb = GAVEUP
    plot2.series[1].format.line.fill.background()
    cat_ax2 = chart2.category_axis
    cat_ax2.tick_labels.font.size = Pt(11); cat_ax2.tick_labels.font.name = SANS
    cat_ax2.tick_labels.font.color.rgb = INK
    cat_ax2.format.line.color.rgb = LINE
    cat_ax2.major_tick_mark = XL_TICK_MARK.NONE
    val_ax2 = chart2.value_axis
    val_ax2.visible = False
    val_ax2.has_major_gridlines = False
    set_chart_font(chart2)

    # stats row
    sx = Inches(0.7)
    line = rounded_box(s, Inches(0.7), Inches(6.15), Inches(11.3), Inches(0.02), fill=LINE, line_color=None, radius=0)
    for num, lab in stats:
        textbox(s, sx, Inches(6.32), Inches(3.3), Inches(0.4), num, font=MONO, size=19, color=INK, bold=True)
        textbox(s, sx, Inches(6.72), Inches(3.3), Inches(0.5), lab, font=MONO, size=9.5, color=FAINT, line_spacing=1.15)
        sx += Inches(3.7)
    return s

s = strategy_chart_slide(
    "13 · Results — everyday tasks", "342 tasks: where the fixes land",
    {"reconstruct": 39, "insert_prep": 37, "local": 15, "already_satisfied": 3, "wrong_action": 2},
    {
        "reconstruct": {"resolved": 39, "gave_up": 0},
        "insert_prep": {"resolved": 37, "gave_up": 0},
        "local": {"resolved": 1, "gave_up": 14},
        "already_satisfied": {"resolved": 3, "gave_up": 0},
        "wrong_action": {"resolved": 2, "gave_up": 0},
    },
    [("96", "FAILURES ACROSS\n68 TASKS"), ("85.4%", "FIXED\nAUTOMATICALLY"), ("1 of 15", "“CAN'T DO THAT”\nPATCHES WORKED")],
)
page_footer(s, 14, TOTAL_SLIDES)

s = strategy_chart_slide(
    "14 · Results — hard tasks", "50 hard tasks: where the fixes land",
    {"insert_prep": 60, "reconstruct": 40, "local": 14, "already_satisfied": 1},
    {
        "reconstruct": {"resolved": 38, "gave_up": 2},
        "insert_prep": {"resolved": 60, "gave_up": 0},
        "already_satisfied": {"resolved": 1, "gave_up": 0},
        "local": {"resolved": 13, "gave_up": 1},
    },
    [("115", "FAILURES ACROSS\n43 TASKS"), ("97.4%", "FIXED\nAUTOMATICALLY"), ("33 of 43", "TASKS USED\nALL 3 TRIES")],
)
page_footer(s, 15, TOTAL_SLIDES)

# =========================================================================
# 16 — Takeaways
# =========================================================================
s = new_slide()
kicker_and_title(s, "15 · Takeaways", "What the numbers actually say", title_size=27)
bullet_list(s, [
    ("The loop fixes 85–97% of what it diagnoses", "From a short list of ideas per failure, every one independently verified.",
     "The remaining gap is mostly one narrow, known weak spot, not broad unreliability."),
    ("The “object can't do that” weak spot depends heavily on the task mix",
     "Only 1 of 15 such patches worked on everyday tasks (driven by “go to sleep”-type tasks), but 13 of 14 worked on the hard set, which doesn't include that task type at all."),
    ("The 3-try limit is a real ceiling, not a formality",
     "33 of 43 repaired hard tasks used all three tries — some recovery headroom is being traded for that tighter limit."),
], Inches(0.7), Inches(2.85), Inches(11.5), item_h=0.68, size=15, gap=0.3)
page_footer(s, 16, TOTAL_SLIDES)

# =========================================================================
# 17 — Closing
# =========================================================================
s = new_slide()
kicker_and_title(s, "Closing", "The short version", title_size=27)
bullet_list(s, [
    ("1", "Most failures aren't random — the system can trace almost all of them back to one exact, identifiable cause."),
    ("2", "A language model supplies judgment about what makes sense; formal rules supply the guarantee that only a genuinely valid fix ever gets used — and that guarantee holds for over 99% of every repair attempted."),
    ("3", "That combination is what turns a plan that breaks constantly into one that recovers on its own, in 9 out of 10 cases, across two very different task sets."),
], Inches(0.7), Inches(2.75), Inches(11.5), item_h=0.85, size=15.5, gap=0.28)
page_footer(s, 17, TOTAL_SLIDES)

# ── save ─────────────────────────────────────────────────────────────────
OUT = "/Users/aytaj/Desktop/embodied-agent-interface-my-version/How_the_Repair_Loop_Works.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
